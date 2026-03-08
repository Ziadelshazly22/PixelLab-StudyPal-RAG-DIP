# -*- coding: utf-8 -*-
"""
app/ui/interface.py
--------------------
Gradio Blocks chat interface for the DIP AI Tutor.

Layout
------
Two tabs:
    1. 💬 Chat   — stateful RAG Q&A via POST /chat
  2. 📄 Upload — PDF ingestion via POST /ingest + status via GET /status

The UI always calls the FastAPI backend at http://localhost:8000.
The optional ``rag_chain`` parameter is kept for backward compatibility
with main.py's ``gr.mount_gradio_app`` call but is not used at runtime.

Standalone launch
-----------------
    python app/ui/interface.py
    # → http://localhost:7860

Mount into FastAPI (main.py)
-----------------------------
    # UI runs as a separate process: python app/ui/interface.py
    import gradio as gr
    from app.ui.interface import build_interface
    app = gr.mount_gradio_app(app, build_interface(), path="/ui")
"""

from __future__ import annotations

import base64
import logging
import re
import uuid

import gradio as gr
from gradio import themes
import requests

# ---------------------------------------------------------------------------
# Favicon — inline base64 data URI so it works regardless of mount path.
# Using a data URI avoids any /favicon.svg path-resolution issues when
# Gradio is mounted at /ui/ inside FastAPI.
# ---------------------------------------------------------------------------
_FAVICON_SVG = (
    '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">'
    '<text y=".9em" font-size="90">\U0001f916</text></svg>'
)
_FAVICON_HEAD = (
    '<link rel="icon" href="data:image/svg+xml;base64,'
    + base64.b64encode(_FAVICON_SVG.encode("utf-8")).decode("ascii")
    + '" type="image/svg+xml">'
)

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Backend configuration
# ---------------------------------------------------------------------------
_BACKEND = "http://localhost:8000"
_CHAT_URL = f"{_BACKEND}/chat"
_SUMMARIZE_URL = f"{_BACKEND}/summarize"
_INGEST_URL = f"{_BACKEND}/ingest"
_STATUS_URL = f"{_BACKEND}/status"
_CHAT_TIMEOUT = 120   # seconds — larger than server's 90s hard limit so server error propagates
_INGEST_TIMEOUT = 60  # seconds — ingestion can be slow


# ---------------------------------------------------------------------------
# Helper: citation formatter
# ---------------------------------------------------------------------------

def _format_citations(text: str) -> str:
    """Bold-emphasize ``[Source: X, Page Y]`` markers in LLM output."""
    return re.sub(
        r"\[Source:\s*([^,\]]+),\s*Page[:\s]*(\d+)\]",
        r"**📖 [Source: \1, Page \2]**",
        text,
    )


# ---------------------------------------------------------------------------
# Helper: call conversational chat API  (stateful, per-session memory)
# ---------------------------------------------------------------------------

def _call_chat_api(
    question: str,
    session_id: str,
    doc_context: str = "",
    doc_filename: str = "",
) -> str:
    """
    POST to ``/chat`` and return the answer.

    When *doc_context* is non-empty the text is sent as a separate JSON field
    so the backend injects it directly into retrieved context — bypassing the
    condense-question step that would otherwise strip the attached content.
    *doc_filename* is forwarded so citations show the real file name.

    Handles:
    - ``ConnectionError`` → backend offline message
    - ``Timeout``         → timeout message
    - HTTP errors         → status code + truncated body
    - Unexpected errors   → generic message with exc string
    """
    # Defensive guard: if Gradio passes the lambda itself (can happen after
    # long idle sessions when State resets to its default value parameter)
    # generate a fresh UUID rather than sending a non-serializable function.
    if callable(session_id) or not isinstance(session_id, str):
        logger.warning("session_id was not a string (%s) — generating new UUID", type(session_id))
        session_id = str(uuid.uuid4())

    payload: dict = {"question": question, "session_id": session_id}
    if doc_context and doc_context.strip():
        payload["doc_context"] = doc_context.strip()
        if doc_filename and doc_filename.strip():
            payload["doc_filename"] = doc_filename.strip()
    try:
        resp = requests.post(_CHAT_URL, json=payload, timeout=_CHAT_TIMEOUT)
        resp.raise_for_status()
        data = resp.json()
        answer = data.get("answer", "")
        sources = data.get("sources", [])
        if sources:
            citations = "\n\n---\n**📚 Sources:**\n"
            for src_item in sources[:3]:
                src = src_item.get("source", "")
                page = src_item.get("page", "")
                if src:
                    citations += (
                        f"- 📖 `{src}`"
                        + (f", Page {page}" if page else "")
                        + "\n"
                    )
            answer = answer + citations
        return _format_citations(answer)

    except requests.exceptions.ConnectionError:
        return (
            "⚠️ Cannot reach the backend server. "
            "Please ensure FastAPI is running on port 8000."
        )
    except requests.exceptions.Timeout:
        return (
            "⚠️ Response timed out. The server may be processing a "
            "large context. Please try again."
        )
    except requests.exceptions.HTTPError as exc:
        body = exc.response.text[:300] if exc.response is not None else str(exc)
        return f"❌ Server error ({exc.response.status_code}): {body}"
    except Exception as exc:
        logger.error("Unexpected error calling /chat API: %s", exc, exc_info=True)
        return f"❌ Unexpected error: {exc}"


# ---------------------------------------------------------------------------
# Helper: call summarize API
# ---------------------------------------------------------------------------

def _call_summarize(
    source: str,
    n_questions: int,
    progress: gr.Progress = gr.Progress(track_tqdm=False),
) -> tuple:
    """
    POST to ``/summarize`` and return ``(summary_markdown, questions_data)``.

    Args:
        source:      Exact filename stored in chunk metadata.
        n_questions: Number of study questions to generate.

    Returns:
        Tuple of ``(summary_str, [[question], ...])`` suitable for
        ``gr.Markdown`` and ``gr.Dataframe`` outputs respectively.
    """
    if not source or not str(source).strip():
        return "⚠️ Please enter a document filename.", []

    payload = {
        "source": str(source).strip(),
        "include_questions": True,
        "n_questions": int(n_questions),
    }
    progress(0.05, desc="Preparing summarize request...")
    try:
        resp = requests.post(
            _SUMMARIZE_URL,
            json=payload,
            timeout=180,  # map-reduce can take 2–3 min for large PDFs
        )
        resp.raise_for_status()
        data = resp.json()
        summary_md = data.get("summary", "⚠️ No summary returned.")
        questions = data.get("study_questions", [])
        questions_data = [[q] for q in questions] if questions else []
        progress(0.75, desc="Received summarize response. Parsing output...")
        progress(1.0, desc="Done")
        return summary_md, questions_data

    except requests.exceptions.ConnectionError:
        return "⚠️ Cannot reach the backend server.", []
    except requests.exceptions.Timeout:
        return "⚠️ Summarisation timed out. The document may be very large.", []
    except requests.exceptions.HTTPError as exc:
        body = exc.response.text[:300] if exc.response is not None else str(exc)
        return f"❌ Server error ({exc.response.status_code}): {body}", []
    except Exception as exc:
        logger.error("Unexpected error calling /summarize API: %s", exc, exc_info=True)
        return f"❌ Unexpected error: {exc}", []


# ---------------------------------------------------------------------------
# Helper: fetch /status
# ---------------------------------------------------------------------------

def _fetch_status() -> str:
    """``GET /status`` → formatted markdown summary string."""
    try:
        resp = requests.get(_STATUS_URL, timeout=5)
        resp.raise_for_status()
        d = resp.json()
        ts = d.get("server_time", "")[:19].replace("T", " ")
        chunks = d.get("total_chunks", "?")
        chunks_fmt = f"{chunks:,}" if isinstance(chunks, int) else str(chunks)
        return (
            f"**Backend:** `{d.get('llm_backend', 'unknown').upper()}`  ·  "
            f"**Embedding:** `{d.get('embedding_model', 'unknown')}`  ·  "
            f"**Chunks:** {chunks_fmt}  ·  "
            f"**Collection:** `{d.get('collection', '?')}`  ·  "
            f"**As of:** {ts} UTC"
        )
    except requests.exceptions.ConnectionError:
        return "⚠️ Backend offline — start the FastAPI server to see live status."
    except Exception as exc:
        return f"⚠️ Could not fetch status: {exc}"


def _fetch_status_and_sources() -> tuple[str, dict]:
    """Return ``(status_markdown, dropdown_update)`` from ``GET /status``."""
    status_md = _fetch_status()
    try:
        resp = requests.get(_STATUS_URL, timeout=5)
        resp.raise_for_status()
        d = resp.json()
        sources = d.get("sources", [])
        source_choices = [s for s in sources if isinstance(s, str) and s.strip()]
        return status_md, gr.update(choices=sorted(set(source_choices)))
    except Exception:
        return status_md, gr.update(choices=[])


# ---------------------------------------------------------------------------
# Helper: upload PDFs
# ---------------------------------------------------------------------------

def _upload_files(files) -> str:
    """
    POST each uploaded PDF to ``/ingest`` as multipart form-data.
    Returns a multi-line status string suitable for a Textbox component.
    """
    if not files:
        return "⚠️ No files selected."

    lines: list[str] = []
    for file_obj in files:
        # Gradio 4.x passes either a NamedString or a file-like with .name
        path = getattr(file_obj, "name", str(file_obj))
        filename = re.split(r"[/\\]", path)[-1]

        try:
            with open(path, "rb") as fh:
                resp = requests.post(
                    _INGEST_URL,
                    files={"file": (filename, fh, "application/pdf")},
                    timeout=_INGEST_TIMEOUT,
                )
            resp.raise_for_status()
            data = resp.json()

            if data.get("status") == "processing":
                lines.append(
                    f"⏳ {filename} — Queued for background ingestion (file > 5 MB). "
                    "Poll 🔄 Refresh Status to check when done."
                )
            else:
                chunks = data.get("chunks_added", 0)
                pages = data.get("pages_processed", 0)
                lines.append(
                    f"✅ {filename} — {chunks:,} chunks added ({pages} pages processed)"
                )

        except requests.exceptions.ConnectionError:
            lines.append(f"❌ {filename} — Error: Cannot reach backend server.")
        except requests.exceptions.HTTPError as exc:
            err = exc.response.text[:200] if exc.response is not None else str(exc)
            lines.append(f"❌ {filename} — Error: {err}")
        except Exception as exc:
            lines.append(f"❌ {filename} — Error: {exc}")

    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Chat event handlers
# ---------------------------------------------------------------------------

# ---------------------------------------------------------------------------
# Helper: extract text from a session-attached document (NOT stored in KB)
# ---------------------------------------------------------------------------

def _extract_session_doc(file_obj) -> tuple[str, str, str]:
    """Extract text from a PDF / DOCX / PPTX for in-session context only.

    Returns ``(extracted_text, status_markdown, filename)``.
    Text is capped at 8 000 chars (~2 000 tokens) to avoid prompt overflow.
    The file is **never** ingested into ChromaDB.
    """
    if file_obj is None:
        return "", "", ""

    path = getattr(file_obj, "name", str(file_obj))
    filename = re.split(r"[/\\]", path)[-1]
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    try:
        text = ""
        if ext == "pdf":
            import fitz  # PyMuPDF — already in requirements
            doc = fitz.open(path)
            text = "\n".join(str(page.get_text()) for page in doc)
            doc.close()
        elif ext in ("docx", "doc"):
            try:
                from docx import Document
                doc = Document(path)
                text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            except ImportError:
                return "", "⚠️ DOCX support requires `python-docx`. Run: `pip install python-docx`", ""
        elif ext in ("pptx", "ppt"):
            try:
                from pptx import Presentation  # type: ignore
                prs = Presentation(path)
                lines = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            shape_text = getattr(shape, "text", "")
                            if shape_text and shape_text.strip():
                                lines.append(shape_text.strip())
                text = "\n".join(lines)
            except ImportError:
                return "", "⚠️ PPTX support requires `python-pptx`. Run: `pip install python-pptx`", ""
        else:
            return "", f"⚠️ Unsupported format `.{ext}`. Use PDF, DOCX, or PPTX.", ""

        if not text.strip():
            return "", f"⚠️ No text found in **{filename}** (image-only or empty).", ""

        MAX_CHARS = 8_000
        truncated = len(text) > MAX_CHARS
        text = text[:MAX_CHARS]
        words = len(text.split())
        note = ", first 8 000 chars shown" if truncated else ""
        status = (
            f"📎 **{filename}** attached ({words:,} words{note}) — "
            f"*session only, not added to knowledge base*"
        )
        return text, status, filename

    except Exception as exc:
        return "", f"❌ Could not read **{filename}**: {exc}", ""


# ---------------------------------------------------------------------------
# Chat event handlers
# ---------------------------------------------------------------------------

def _handle_send(
    user_message: str,
    chat_history: list,
    session_id: str,
    doc_context: str,
    doc_filename: str = "",
):
    """Stream user message + RAG answer to the chatbot.

    Yields an immediate '⏳ Thinking…' placeholder so the user gets visual
    feedback within < 1 s, then replaces it with the real answer once the
    backend returns (or an error message on timeout / quota exhaustion).
    """
    if not user_message.strip():
        yield chat_history, user_message
        return

    # 1. Instant feedback — show the question + placeholder right away
    yield (
        chat_history + [
            {"role": "user", "content": user_message},
            {"role": "assistant", "content": "⏳ *Searching knowledge base… (up to 30 s)*"},
        ],
        "",
    )

    # 2. Call the backend (blocks this thread; Gradio runs it in a worker)
    answer = _call_chat_api(user_message.strip(), session_id, doc_context, doc_filename)

    # 3. Replace placeholder with the real answer
    yield (
        chat_history + [
            {"role": "user", "content": user_message},
            {"role": "assistant", "content": answer},
        ],
        "",
    )


def _handle_clear(session_id: str) -> tuple[list, str]:
    """Wipe chat history locally and clear server-side session memory."""
    try:
        requests.delete(f"{_BACKEND}/chat/{session_id}", timeout=5)
    except Exception:
        pass  # Best-effort — don't block the UI clear if backend is unreachable
    return [], ""


# ---------------------------------------------------------------------------
# Interface builder
# ---------------------------------------------------------------------------

def build_interface(rag_chain=None) -> gr.Blocks:
    """
    Construct and return the full Gradio Blocks interface.

    Parameters
    ----------
    rag_chain : optional
        Accepted but unused — kept for compatibility with
        ``gr.mount_gradio_app`` calls in main.py.

    Returns
    -------
    gr.Blocks
    """
    with gr.Blocks(
        title="DIP AI Tutor",
        # Inline base64 favicon — works regardless of mount path (/ui/ vs /).
        # A bare href="/favicon.svg" is resolved relative to the document base
        # and can be shadowed by the browser's auto-discovery; a data URI is
        # self-contained and always wins.
        head=_FAVICON_HEAD,
    ) as demo:

        # ── Session state — set via demo.load() so Gradio never passes the
        # lambda itself to event handlers (which happens after long idle
        # sessions when the State resets to its default value parameter).
        session_id = gr.State(value=None)

        # ── Top status bar (always visible) ─────────────────────────────
        with gr.Row():
            status_bar = gr.Markdown(
                value="⏳ Connecting to backend...",
                elem_id="status-bar",
            )

        # ── Main tabs ───────────────────────────────────────────────────
        with gr.Tabs():

            # ── TAB 1: CHAT ─────────────────────────────────────────────
            with gr.Tab("💬 Chat"):
                gr.Markdown(
                    "# 🎓 Digital Image Processing AI Tutor\n"
                    "Powered by **Gonzalez & Woods 4th Ed** · OpenCV · NumPy · SciPy"
                )

                # ── Session doc states (not in KB, cleared on page reload) ──
                doc_context  = gr.State(value="")   # extracted text
                doc_filename = gr.State(value="")   # original filename for citations

                # ── Session-only document attachment (always visible) ────────
                with gr.Group(elem_id="attach-group"):
                    gr.Markdown(
                        "### 📎 Attach a Document *(Session Only)*\n"
                        "Upload a **PDF, DOCX, or PPTX** to discuss alongside your DIP questions. "
                        "⚠️ *This file is never added to the knowledge base.*"
                    )
                    with gr.Row(equal_height=True):
                        session_file = gr.File(
                            file_types=[".pdf", ".docx", ".doc", ".pptx", ".ppt"],
                            file_count="single",
                            label="Select document (PDF / DOCX / PPTX)",
                            scale=5,
                        )
                        with gr.Column(scale=1, min_width=140):
                            attach_btn = gr.Button("📎 Attach", variant="secondary", size="lg")
                            detach_btn = gr.Button("✖ Remove", variant="stop", size="sm")
                    attach_info = gr.Markdown(
                        value="*No document attached — chat answers from the DIP knowledge base only.*",
                        elem_id="attach-info",
                    )

                attach_btn.click(
                    fn=_extract_session_doc,
                    inputs=[session_file],
                    outputs=[doc_context, attach_info, doc_filename],
                )
                detach_btn.click(
                    fn=lambda: ("", "*No document attached — chat answers from the DIP knowledge base only.*", None, ""),
                    inputs=[],
                    outputs=[doc_context, attach_info, session_file, doc_filename],
                )

                gr.Markdown("---")

                chatbot = gr.Chatbot(
                    height=460,
                    show_label=False,
                    elem_id="chatbot",
                    latex_delimiters=[
                        {"left": "$$", "right": "$$", "display": True},   # display math
                        {"left": "$",  "right": "$",  "display": False},  # inline math
                        {"left": "\\[", "right": "\\]", "display": True}, # \[...\]
                        {"left": "\\(", "right": "\\)", "display": False}, # \(...\)
                    ],
                )

                question_box = gr.Textbox(
                    placeholder=(
                        "Ask about spatial filtering, Fourier transforms, "
                        "morphological operations..."
                    ),
                    label="Your Question",
                    lines=2,
                )

                with gr.Row():
                    send_btn = gr.Button("Send", variant="primary")
                    clear_btn = gr.Button("🗑️ Clear Conversation")

                with gr.Accordion("💡 Example Questions", open=False):
                    _EXAMPLES = [
                        "What is histogram equalization and when is it used?",
                        "Derive the discrete Fourier Transform for 2D images.",
                        "Explain morphological erosion vs dilation with OpenCV code.",
                        "What noise models are common in DIP and how do we remove them?",
                        "How does the Canny edge detector work step by step?",
                    ]
                    for _ex in _EXAMPLES:
                        _btn = gr.Button(
                            _ex, size="sm", elem_classes=["example-btn"]
                        )
                        # Each button click populates the textbox
                        _btn.click(
                            fn=lambda e=_ex: e,
                            inputs=[],
                            outputs=[question_box],
                        )

                # ── Chat events ─────────────────────────────────────────
                send_btn.click(
                    fn=_handle_send,
                    inputs=[question_box, chatbot, session_id, doc_context, doc_filename],
                    outputs=[chatbot, question_box],
                )
                question_box.submit(
                    fn=_handle_send,
                    inputs=[question_box, chatbot, session_id, doc_context, doc_filename],
                    outputs=[chatbot, question_box],
                )
                clear_btn.click(
                    fn=_handle_clear,
                    inputs=[session_id],
                    outputs=[chatbot, question_box],
                )

            # ── TAB 2: UPLOAD ────────────────────────────────────────────
            with gr.Tab("📄 Upload Documents"):

                # ── Section header ──────────────────────────────────────
                gr.Markdown(
                    "## 📥 Add Documents to Knowledge Base\n"
                    "Upload one or more **PDF** files to expand the tutor's "
                    "knowledge base. Each file is chunked, embedded, and "
                    "stored in ChromaDB immediately."
                )

                # ── Upload row: file picker + button side by side ───────
                with gr.Row(equal_height=True):
                    with gr.Column(scale=4):
                        file_upload = gr.File(
                            file_types=[".pdf"],
                            file_count="multiple",
                            label="📂 Select PDF(s) to upload",
                            height=160,
                        )
                    with gr.Column(scale=1, min_width=180):
                        upload_btn = gr.Button(
                            "⬆️ Add to Knowledge Base",
                            variant="primary",
                            size="lg",
                        )
                        refresh_btn = gr.Button(
                            "🔄 Refresh Status",
                            variant="secondary",
                            size="sm",
                        )

                # ── Ingestion result ────────────────────────────────────
                ingestion_status = gr.Textbox(
                    label="📋 Ingestion Result",
                    placeholder="Upload a PDF above — results will appear here.",
                    interactive=False,
                    lines=4,
                    max_lines=10,
                )

                # ── Knowledge-base status panel ─────────────────────────
                with gr.Accordion("📊 Knowledge-Base Status", open=True):
                    status_display = gr.Markdown(
                        value="⏳ Loading status…",
                        elem_id="status_display",
                    )

                gr.Markdown("---")

                # ── Summarize section ───────────────────────────────────
                gr.Markdown(
                    "## 🔍 Summarize a Document\n"
                    "Generate an academic summary and exam-style study "
                    "questions for any ingested PDF. "
                    "_This may take 1–3 minutes for large documents._"
                )

                with gr.Row(equal_height=True):
                    with gr.Column(scale=3):
                        summarize_filename = gr.Dropdown(
                            choices=[],
                            allow_custom_value=True,
                            label="📄 Document filename",
                            info="Select from the knowledge base or type a filename manually.",
                        )
                    with gr.Column(scale=1, min_width=200):
                        n_questions_slider = gr.Slider(
                            minimum=2,
                            maximum=10,
                            value=5,
                            step=1,
                            label="🎯 Number of Study Questions",
                        )

                summarize_btn = gr.Button(
                    "📝 Generate Summary & Study Questions",
                    variant="primary",
                    size="lg",
                )

                summarize_status = gr.Markdown(
                    value="",
                    visible=True,
                    elem_id="summarize-status",
                )

                with gr.Accordion("📄 Document Summary", open=True):
                    summary_output = gr.Markdown(
                        value="_Summary will appear here after generation._",
                    )

                with gr.Accordion("🎓 Study Questions", open=True):
                    questions_output = gr.Dataframe(
                        headers=["Study Questions"],
                        label="🎓 Study Questions",
                        interactive=False,
                        wrap=True,
                    )

                # ── Upload / status events ──────────────────────────────
                upload_btn.click(
                    fn=_upload_files,
                    inputs=[file_upload],
                    outputs=[ingestion_status],
                ).then(
                    fn=_fetch_status_and_sources,
                    inputs=[],
                    outputs=[status_display, summarize_filename],
                )

                refresh_btn.click(
                    fn=_fetch_status_and_sources,
                    inputs=[],
                    outputs=[status_display, summarize_filename],
                )

                # ── Summarize events ────────────────────────────────────
                summarize_btn.click(
                    fn=lambda: (
                        "⏳ Generating summary… this may take 1–3 minutes. Please wait.",
                        "_Generating…_",
                        [],
                    ),
                    inputs=[],
                    outputs=[summarize_status, summary_output, questions_output],
                    queue=False,
                ).then(
                    fn=_call_summarize,
                    inputs=[summarize_filename, n_questions_slider],
                    outputs=[summary_output, questions_output],
                ).then(
                    fn=lambda: "",
                    inputs=[],
                    outputs=[summarize_status],
                )
        # ── On page load: (1) populate status bar, (2) generate session UUID ──
        # Using demo.load() for session_id guarantees Gradio always receives a
        # proper string UUID — never the lambda that was set as the default value
        # (which Gradio 6 can pass as-is after long idle reconnections).
        demo.load(
            fn=_fetch_status_and_sources,
            inputs=[],
            outputs=[status_bar, summarize_filename],
        )
        demo.load(
            fn=lambda: str(uuid.uuid4()),
            inputs=[],
            outputs=[session_id],
        )

    return demo


# ---------------------------------------------------------------------------
# Standalone launch
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    logging.basicConfig(
        level=logging.INFO,
        format="%(asctime)s | %(levelname)s | %(message)s",
        datefmt="%H:%M:%S",
    )
    import gradio.utils as _gu
    _gu.get_favicon_path = lambda: None  # type: ignore
    demo = build_interface()
    demo.launch(
        server_name="0.0.0.0",
        server_port=7860,
        share=False,
        show_error=True,
        favicon_path="app/ui/favicon.svg",
        theme=themes.Soft(),
        css="#status-bar{font-size:.80em;opacity:.88;padding:4px 0} .example-btn{margin:2px 0!important}",
    )
